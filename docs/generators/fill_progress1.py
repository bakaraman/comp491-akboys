from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation("lecture_slides/Comp491_ProgressPresentationTemplate_fixed.pptx")


def clear_and_set(shape, lines, font_size=18, bold_first=False):
    """Replace all paragraphs in a text frame with new lines."""
    tf = shape.text_frame
    # Clear existing paragraphs
    for i in range(len(tf.paragraphs) - 1, 0, -1):
        p = tf.paragraphs[i]._p
        p.getparent().remove(p)
    # First line uses existing paragraph
    for li, line in enumerate(lines):
        if li == 0:
            para = tf.paragraphs[0]
            para.clear()
        else:
            para = tf.add_paragraph()

        # Handle tuple (text, size, bold, color)
        if isinstance(line, tuple):
            text, sz, bld, color = line
        else:
            text = line
            sz = font_size
            bld = (li == 0 and bold_first)
            color = None

        run = para.add_run()
        run.text = text
        run.font.size = Pt(sz)
        run.font.bold = bld
        if color:
            run.font.color.rgb = RGBColor(*color)
        para.space_after = Pt(3)
        para.space_before = Pt(0)


# ── Slide 0: Title Slide ───────────────────────────────────────
slide0 = prs.slides[0]
txBox = slide0.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(8), Inches(2.5))
tf = txBox.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.space_after = Pt(12)
r = p.add_run()
r.text = "Text-Based Adventure with LLMs"
r.font.size = Pt(32)
r.font.bold = True
r.font.color.rgb = RGBColor(255, 255, 255)

p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
p2.space_after = Pt(8)
r2 = p2.add_run()
r2.text = "COMP 491 — Progress Meeting 1"
r2.font.size = Pt(20)
r2.font.color.rgb = RGBColor(220, 220, 220)

p3 = tf.add_paragraph()
p3.alignment = PP_ALIGN.CENTER
r3 = p3.add_run()
r3.text = "March 5, 2026"
r3.font.size = Pt(16)
r3.font.color.rgb = RGBColor(200, 200, 200)


# ── Slide 1: Project Info ──────────────────────────────────────
slide1 = prs.slides[1]
for shape in slide1.shapes:
    if shape.name == "Content Placeholder 2":
        clear_and_set(shape, [
            ("Project Title: Text-Based Adventure with LLMs", 18, True, None),
            "",
            ("Team Members:", 16, True, None),
            ("  Kadir Yigit Ozcelik (79975)", 15, False, None),
            ("  Serdar Yengil (80232)", 15, False, None),
            ("  Batuhan Karaman (79791)", 15, False, None),
            ("  Ata Berke Goktekin (80277)", 15, False, None),
            "",
            ("Advisor: Baris Akgun", 16, True, None),
            ("Progress Meeting 1 — March 5, 2026", 15, False, None),
        ], font_size=15)


# ── Slide 2: Previous Progress Meeting ─────────────────────────
slide2 = prs.slides[2]
for shape in slide2.shapes:
    if shape.name == "Content Placeholder 2":
        clear_and_set(shape, [
            ("This is the first progress meeting.", 16, True, None),
            "",
            ("Completed before this meeting:", 16, True, None),
            ("• Met with advisor (Baris Akgun) to discuss project scope", 14, False, None),
            ("• Advisor emphasized multiplayer as a key differentiator", 14, False, None),
            ("• Built a working single-player demo (Tkinter, Gemini API, 5 rooms, 3 NPCs)", 14, False, None),
            ("• Prepared and submitted Project Registration Form", 14, False, None),
            ("• Prepared and submitted Project Proposal with Gantt chart", 14, False, None),
            ("• Decided on tech stack: GCP (Vertex AI, Cloud Run, Firebase)", 14, False, None),
        ], font_size=14)


# ── Slide 3: Current Progress Meeting ──────────────────────────
slide3 = prs.slides[3]
for shape in slide3.shapes:
    if shape.name == "Content Placeholder 2":
        clear_and_set(shape, [
            ("WP1 — Research, Setup & Demo Refinement (W1-W2)", 15, True, None),
            "",
            ("Tasks Performed:", 14, True, None),
            ("• GitHub repo created with branch protection and issue tracking", 13, False, None),
            ("• Demo refactored: game engine separated from LLM layer", 13, False, None),
            ("• Architecture documented (Structured Chaos: Python referee + LLM narrator)", 13, False, None),
            ("• Project proposal and registration form prepared and submitted", 13, False, None),
            "",
            ("Work Distribution:", 14, True, None),
            ("• Batuhan: Demo development, architecture design, proposal", 13, False, None),
            ("• Kadir: Technical document, function call design", 13, False, None),
            ("• Serdar: Registration form, multiplayer research", 13, False, None),
            ("• Ata Berke: Background research, references", 13, False, None),
            "",
            ("Deliverables:", 14, True, None),
            ("• D1.1 GitHub repo with CI  • D1.2 Working demo  • D1.3 Proposal + Reg. Form", 13, False, None),
        ], font_size=13)


# ── Slide 4: Next Progress Meeting ─────────────────────────────
slide4 = prs.slides[4]
for shape in slide4.shapes:
    if shape.name == "Content Placeholder 2":
        clear_and_set(shape, [
            ("Next: WP2 — Procedural Story Generation & Core Engine (W3-W6)", 15, True, None),
            "",
            ("Planned Tasks:", 14, True, None),
            ("• Design scenario JSON schema and validation rules", 13, False, None),
            ("• Implement Phase 1 LLM scenario generation with re-prompt loop", 13, False, None),
            ("• Build schema validator (graph connectivity, evidence chain)", 13, False, None),
            ("• Start 5 universe templates (noir, cyberpunk, fantasy, horror, western)", 13, False, None),
            ("• Begin FastAPI backend migration (from Tkinter demo)", 13, False, None),
            "",
            ("Work Distribution:", 14, True, None),
            ("• Batuhan: JSON schema, Phase 1 LLM generation, FastAPI backend", 13, False, None),
            ("• Ata Berke: Schema validator, universe templates", 13, False, None),
            ("• Kadir: React/Next.js frontend setup (WP3 prep)", 13, False, None),
            ("• Serdar: WebSocket/multiplayer research (WP4 prep)", 13, False, None),
        ], font_size=13)


# ── Slide 5: References ────────────────────────────────────────
slide5 = prs.slides[5]
for shape in slide5.shapes:
    if shape.name == "Content Placeholder 2":
        clear_and_set(shape, [
            ("[1] N. Montfort, Twisty Little Passages: An Approach to Interactive Fiction. MIT Press, 2003.",
             12, False, None),
            ("[2] A. Plotkin, The Inform Designer's Handbook, 2004.",
             12, False, None),
            ("[3] P. Ammanabrolu et al., How to Avoid Being Eaten by a Grue, arXiv:2006.07409, 2020.",
             12, False, None),
            ("[4] N. Walton, AI Dungeon, Latitude, 2019. https://aidungeon.io",
             12, False, None),
            ("[5] Google DeepMind, Gemini API: Function Calling, 2024. https://ai.google.dev/docs/function_calling",
             12, False, None),
            ("[6] N. Shaker et al., Procedural Content Generation in Games. Springer, 2016.",
             12, False, None),
        ])


# ── Slide 6: Attachment ────────────────────────────────────────
slide6 = prs.slides[6]
for shape in slide6.shapes:
    if shape.name == "Content Placeholder 2":
        clear_and_set(shape, [
            ("Submitted Documents:", 16, True, None),
            ("• COMP491_Registration_Form_FILLED.docx", 14, False, None),
            ("• COMP491_Proposal_FILLED.docx (with Gantt chart)", 14, False, None),
            "",
            ("Working Demo:", 16, True, None),
            ("• Single-player noir detective game (main.py)", 14, False, None),
            ("• Gemini 2.5 Flash API, Tkinter GUI, 6 function calls", 14, False, None),
            ("• GitHub repository: (link to be added)", 14, False, None),
        ], font_size=14)


# ── Save ───────────────────────────────────────────────────────
out = "/Users/batuhankaraman/demo491/COMP491_Progress1.pptx"
prs.save(out)
print(f"Saved: {out}")
