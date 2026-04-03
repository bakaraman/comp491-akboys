from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation("../lecture-slides/Comp491_ProgressPresentationTemplate_fixed.pptx")


def clear_and_set(shape, lines, font_size=18, bold_first=False):
    """Replace all paragraphs in a text frame with new lines."""
    tf = shape.text_frame
    for i in range(len(tf.paragraphs) - 1, 0, -1):
        p = tf.paragraphs[i]._p
        p.getparent().remove(p)
    for li, line in enumerate(lines):
        if li == 0:
            para = tf.paragraphs[0]
            para.clear()
        else:
            para = tf.add_paragraph()

        if isinstance(line, tuple):
            text, sz, bld, color = line
        else:
            text = line
            sz = font_size
            bld = (li == 0 and bold_first)
            color = None

        run = para.add_run()
        run.text = text
        run.font.size = Pt(sz)
        run.font.bold = bld
        if color:
            run.font.color.rgb = RGBColor(*color)
        para.space_after = Pt(3)
        para.space_before = Pt(0)


# ── Slide 0: Title Slide ───────────────────────────────────────
slide0 = prs.slides[0]
txBox = slide0.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(8), Inches(2.5))
tf = txBox.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.space_after = Pt(12)
r = p.add_run()
r.text = "Text-Based Adventure with LLMs"
r.font.size = Pt(32)
r.font.bold = True
r.font.color.rgb = RGBColor(255, 255, 255)

p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
p2.space_after = Pt(8)
r2 = p2.add_run()
r2.text = "COMP 491 — Progress Meeting 3"
r2.font.size = Pt(20)
r2.font.color.rgb = RGBColor(220, 220, 220)

p3 = tf.add_paragraph()
p3.alignment = PP_ALIGN.CENTER
r3 = p3.add_run()
r3.text = "April 3, 2026"
r3.font.size = Pt(16)
r3.font.color.rgb = RGBColor(200, 200, 200)


# ── Slide 1: Project Info ──────────────────────────────────────
slide1 = prs.slides[1]
for shape in slide1.shapes:
    if shape.name == "Content Placeholder 2":
        clear_and_set(shape, [
            ("Project Title: Text-Based Adventure with LLMs", 18, True, None),
            "",
            ("Team Members:", 16, True, None),
            ("  Kadir Yigit Ozcelik (79975)", 15, False, None),
            ("  Serdar Yengil (80232)", 15, False, None),
            ("  Batuhan Karaman (79791)", 15, False, None),
            ("  Ata Berke Goktekin (80277)", 15, False, None),
            "",
            ("Advisor: Baris Akgun", 16, True, None),
            ("Progress Meeting 3 — April 3, 2026", 15, False, None),
        ], font_size=15)


# ── Slide 2: Previous Progress Meeting ─────────────────────────
slide2 = prs.slides[2]
for shape in slide2.shapes:
    if shape.name == "Content Placeholder 2":
        clear_and_set(shape, [
            ("Progress Meeting 2 Summary (March 12, 2026):", 16, True, None),
            "",
            ("What was accomplished in WP2:", 14, True, None),
            ("• Full migration from Python/Tkinter to TypeScript monorepo", 12, False, None),
            ("• Next.js 15 + React 19 frontend, Express.js backend", 12, False, None),
            ("• 6 game scenarios (Noir, Haunted, Space, Pirate, Western, Cyberpunk)", 12, False, None),
            ("• OpenAI GPT-5.4 streaming narration via SSE", 12, False, None),
            ("• GPT-5-nano follow-up action suggestions", 12, False, None),
            ("• Session management with UUID routing, markdown rendering", 12, False, None),
            "",
            ("What was planned for WP3:", 14, True, None),
            ("• Socket.IO multiplayer co-op", 12, False, None),
            ("• Firestore session persistence", 12, False, None),
            ("• Firebase Auth (Google sign-in)", 12, False, None),
            ("• Cloud Run + Firebase Hosting deployment", 12, False, None),
            ("• AI image generation (gpt-image-1.5)", 12, False, None),
            ("• Game-over conditions and accusation mechanic", 12, False, None),
        ], font_size=12)


# ── Slide 3: Current Progress Meeting ──────────────────────────
slide3 = prs.slides[3]
for shape in slide3.shapes:
    if shape.name == "Content Placeholder 2":
        clear_and_set(shape, [
            ("WP3+WP4 — Multiplayer, Persistence, Deploy & Advanced Features", 14, True, None),
            ("All 8 planned issues completed, 6 PRs merged.", 12, True, (0, 128, 0)),
            "",
            ("Kadir (#3, #5) — Multiplayer System:", 13, True, None),
            ("• Socket.IO real-time engine with room codes and lobby system", 11, False, None),
            ("• Multiplayer UI: player sidebar, typing indicators, comm panel", 11, False, None),
            ("• Action batching with configurable timing windows", 11, False, None),
            ("• Scoped visibility (private vs observed messages per player)", 11, False, None),
            "",
            ("Serdar (#4, #6) — Auth & GameState:", 13, True, None),
            ("• Firebase Auth middleware with Google sign-in", 11, False, None),
            ("• Server-side GameState tracking (room, inventory, evidence, turns)", 11, False, None),
            ("• Structured directive validation (MOVE, PICKUP, UNLOCK, etc.)", 11, False, None),
            "",
            ("Ata Berke (#7, #8) — Visuals & Game Logic:", 13, True, None),
            ("• gpt-image-1.5 scene generation with per-scenario style prompts", 11, False, None),
            ("• Image caching per room (no redundant API calls)", 11, False, None),
            ("• Accusation mechanic with evidence chain validation", 11, False, None),
            ("• Win/lose states and replay reset behavior", 11, False, None),
            "",
            ("Batuhan (#1, #2) — Infrastructure:", 13, True, None),
            ("• Backend deployed to Google Cloud Run (auto-scaling)", 11, False, None),
            ("• Frontend deployed to Firebase Hosting (velvet-shadow.web.app)", 11, False, None),
            ("• Firestore-backed session persistence (survives restarts)", 11, False, None),
            ("• Production CORS, env config, Firebase project setup", 11, False, None),
        ], font_size=11)


# ── Slides 4-7: Individual member pages (right after current progress) ──
# We use the template's slide 4 for overview, then add 4 new slides for members.

# First fill the original slide 4 as Kadir's page
slide4 = prs.slides[4]
for shape in slide4.shapes:
    if shape.name == "Content Placeholder 2":
        clear_and_set(shape, [
            ("Completed (Issues #3, #5):", 13, True, None),
            ("• Built the entire multiplayer frontend from scratch, because the single-player chat UI", 10, False, None),
            ("  could not handle multiple players — a lobby, voting screen, and shared game view were needed", 10, False, None),
            ("• Integrated Socket.IO on the client side so the browser keeps a persistent connection to the", 10, False, None),
            ("  server, which allows real-time updates without polling — messages appear instantly for everyone", 10, False, None),
            ("• Added a slide-out player sidebar that shows each player's name, color, online status, current", 10, False, None),
            ("  room, and inventory, because in multiplayer you need to know where your teammates are", 10, False, None),
            ("• Built the communication panel as a separate slide-out drawer with two tabs: room chat (everyone", 10, False, None),
            ("  in your room hears it) and direct messages (private to one player), because players need to", 10, False, None),
            ("  coordinate without the narrator seeing their conversation", 10, False, None),
            ("• Implemented scoped message visibility — when the narrator responds to Player A, only Player A", 10, False, None),
            ("  sees the full response; other players in the same room get a short third-person summary instead", 10, False, None),
            "",
            ("Next Up (7 issues):", 13, True, None),
            ("• Will add push-to-talk voice chat (#15) using WebRTC, because typing breaks immersion in a", 10, False, None),
            ("  roleplay game — holding Space will broadcast to your room, holding T activates the walkie-talkie", 10, False, None),
            ("• Will build an interactive map (#21) that renders rooms as nodes with fog of war, because players", 10, False, None),
            ("  currently have no visual sense of the game world — unvisited rooms stay hidden until explored", 10, False, None),
            ("• Will create an evidence board (#22) where players pin collected clues and draw connections", 10, False, None),
            ("  between suspects, because right now all evidence is buried in the chat scroll", 10, False, None),
            ("• Will add a multiplayer game-over screen (#23), because currently multiplayer games never end —", 10, False, None),
            ("  there is no win/lose overlay, no accuse button, and no turn counter in the multiplayer UI", 10, False, None),
            ("• Will add ambient audio and sound effects (#24) per scenario — rain and jazz for noir, creaking", 10, False, None),
            ("  wood for haunted manor — because the game is completely silent right now", 10, False, None),
            ("• Will build a post-game reconstruction (#36) that replays what actually happened step by step", 10, False, None),
            ("• Will design an NPC interrogation panel (#37) with a trust meter, because talking to NPCs", 10, False, None),
            ("  currently feels the same as typing any other action", 10, False, None),
        ], font_size=11)
    elif shape.name == "Title 1":
        shape.text_frame.paragraphs[0].clear()
        r = shape.text_frame.paragraphs[0].add_run()
        r.text = "Kadir Yigit Ozcelik — Frontend & UI"
        r.font.size = Pt(28)
        r.font.bold = True

# Now add 3 more slides for other members using add_slide
from copy import deepcopy

def add_member_slide(prs, title_text, content_lines):
    layout = prs.slides[4].slide_layout
    slide = prs.slides.add_slide(layout)
    for shape in slide.shapes:
        if shape.name == "Title 1":
            shape.text_frame.paragraphs[0].clear()
            r = shape.text_frame.paragraphs[0].add_run()
            r.text = title_text
            r.font.size = Pt(28)
            r.font.bold = True
        elif shape.name == "Content Placeholder 2":
            clear_and_set(shape, content_lines, font_size=11)
    return slide

# Serdar's page
add_member_slide(prs, "Serdar Yengil — Backend & API", [
    ("Completed (Issues #4, #6):", 13, True, None),
    ("• Added Firebase Auth middleware to the backend, because previously anyone could call the API without", 10, False, None),
    ("  logging in — now every request must include a valid Google sign-in token or it gets rejected", 10, False, None),
    ("• Built server-side GameState tracking, because before this the AI was the only thing that knew which", 10, False, None),
    ("  room the player was in or what items they had — now the server tracks room, inventory, visited rooms,", 10, False, None),
    ("  discovered evidence, and turn count independently from the AI", 10, False, None),
    ("• Implemented structured directive validation — the narrator returns JSON directives like MOVE, PICKUP,", 10, False, None),
    ("  UNLOCK, and the server checks each one against the actual game state before applying it, so the AI", 10, False, None),
    ("  cannot teleport players to nonexistent rooms or give them items that don't exist", 10, False, None),
    ("• Created the login page with Google sign-in and protected all game routes behind authentication", 10, False, None),
    "",
    ("Next Up (8 issues):", 13, True, None),
    ("• Will add multiplayer game-over (#25), because right now multiplayer games never end — there is no", 10, False, (200, 0, 0)),
    ("  turn counting, no accusation mechanic, and no win/lose state; this is the most critical missing piece", 10, False, (200, 0, 0)),
    ("• Will implement player roles (#18) — Detective can accuse, Thief finds hidden rooms, Doctor resists", 10, False, None),
    ("  sanity damage, Journalist gets bonus NPC dialogue — because currently all players are identical", 10, False, None),
    ("• Will add NPC cross-player memory (#19), because if Player A threatens an NPC, that NPC should be", 10, False, None),
    ("  hostile to Player B too — right now NPCs treat each player as a stranger every time", 10, False, None),
    ("• Will build evidence chain prerequisites (#26), because currently all clues can be found in any order", 10, False, None),
    ("  — in a real detective game, finding the poison signs on the body should be required before finding", 10, False, None),
    ("  the hidden poison bottle in the garden", 10, False, None),
    ("• Will replace the text-matching evidence discovery (#27) with structured LLM extraction, because the", 10, False, None),
    ("  current system falsely discovers evidence when the AI mentions an item name in passing", 10, False, None),
    ("• Will bring back the sanity system (#38) from the original demo — dangerous actions drain sanity and", 10, False, None),
    ("  hitting zero means game over, which adds real risk to exploration", 10, False, None),
    ("• Will make NPCs move between rooms (#39), because they currently stand in one spot forever", 10, False, None),
    ("• Will add hidden rooms (#40) discoverable through investigation or items, adding exploration depth", 10, False, None),
])

# Ata Berke's page
add_member_slide(prs, "Ata Berke Goktekin — Game Design & Prompts", [
    ("Completed (Issues #7, #8):", 13, True, None),
    ("• Built the AI image generation pipeline using gpt-image-1.5 — every time a player enters a new room,", 10, False, None),
    ("  the server generates a scene image using a style prompt matched to the scenario's theme, so noir gets", 10, False, None),
    ("  '1920s ink illustration' and cyberpunk gets 'neon-drenched concept art'", 10, False, None),
    ("• Added per-room image caching, because without it every room revisit would trigger another expensive", 10, False, None),
    ("  API call — now the image is generated once and served from cache on subsequent visits", 10, False, None),
    ("• Implemented the accusation mechanic — the player picks a suspect and a piece of evidence, and the", 10, False, None),
    ("  server checks three conditions: correct suspect, correct key evidence, and all required evidence", 10, False, None),
    ("  items discovered; if all three match, the player wins, otherwise the game is lost", 10, False, None),
    ("• Added the game-over overlay with win (gold, 'Mystery Solved') and lose (red, 'Game Over') screens,", 10, False, None),
    ("  plus a replay button that creates a fresh session with the same scenario", 10, False, None),
    "",
    ("Next Up (7 issues):", 13, True, None),
    ("• Will write three new scenarios (#28) — Medieval Castle (a king poisoned at a feast), Arctic Research", 10, False, None),
    ("  Station (a scientist vanishes during a blizzard), and Underwater Base (sabotage at the ocean floor)", 10, False, None),
    ("• Will add a dynamic hint system (#29), because players often get stuck and waste all 12 turns in the", 10, False, None),
    ("  wrong room — the narrator will start dropping subtle hints at turn 5 and obvious ones by turn 11", 10, False, None),
    ("• Will introduce difficulty levels (#30), because all scenarios currently have the same 12-turn limit —", 10, False, None),
    ("  easy scenarios will get 15 turns and hard ones will get 9, shown as a badge on the scenario picker", 10, False, None),
    ("• Will fix known scenario bugs (#31) — the haunted manor cellar is defined but has no entrance from the", 10, False, None),
    ("  foyer, and noir requires 4 evidence items while every other scenario requires only 3", 10, False, None),
    ("• Will add red herrings (#41) — false evidence items that look real but lead to wrong accusations,", 10, False, None),
    ("  because right now every clue is genuine and the mystery is too straightforward", 10, False, None),
    ("• Will build an NPC alibi system (#42) so players can cross-reference what different NPCs claim and", 10, False, None),
    ("  spot contradictions in the guilty NPC's story", 10, False, None),
    ("• Will implement procedural scenario generation (#43) where the LLM creates an entirely new mystery", 10, False, None),
    ("  from a theme prompt, validated by the server for solvability — this enables infinite replayability", 10, False, None),
])

# Batuhan's page
add_member_slide(prs, "Batuhan Karaman — Infrastructure & GCP", [
    ("Completed (Issues #1, #2):", 13, True, None),
    ("• Deployed the backend to Google Cloud Run in europe-west1, because the app was previously only", 10, False, None),
    ("  accessible on localhost — Cloud Run auto-scales based on traffic and costs nothing when idle", 10, False, None),
    ("• Deployed the frontend to Firebase Hosting at velvet-shadow.web.app with a global CDN, so the app", 10, False, None),
    ("  loads fast from anywhere in the world without requiring users to install anything", 10, False, None),
    ("• Migrated the session store from an in-memory Map to Firestore, because with the old approach every", 10, False, None),
    ("  server restart or Cloud Run scaling event wiped all active game sessions — now sessions are persisted", 10, False, None),
    ("  in Firestore and automatically rehydrated when the server starts", 10, False, None),
    ("• Configured production CORS so the frontend (velvet-shadow.web.app) and backend (Cloud Run) can", 10, False, None),
    ("  communicate across different domains without the browser blocking requests", 10, False, None),
    ("• Set up the Firebase project with Google Auth provider, Firestore database, and all environment", 10, False, None),
    ("  variables needed for production — separate configs for local dev vs deployed environments", 10, False, None),
    "",
    ("Next Up (5 issues):", 13, True, None),
    ("• Will set up a CI/CD pipeline (#32) with GitHub Actions, because right now every deploy is done", 10, False, None),
    ("  manually from the terminal — pushing to main will auto-build, run type checks, and deploy to both", 10, False, None),
    ("  Cloud Run and Firebase Hosting", 10, False, None),
    ("• Will move scene images to Cloud Storage (#33), because they are currently cached in server memory —", 10, False, None),
    ("  every restart loses them, and different Cloud Run instances don't share the cache", 10, False, None),
    ("• Will add rate limiting (#34), because there is currently nothing stopping someone from spamming the", 10, False, None),
    ("  API and running up the OpenAI bill — per-IP, per-session, and daily cost caps will be enforced", 10, False, None),
    ("• Will build an analytics dashboard (#35) to track how many people are playing, which scenarios are", 10, False, None),
    ("  popular, how much the API costs per day, and what the win/lose rate is per scenario", 10, False, None),
    ("• Will add spectator and replay modes (#44) so people can watch live games or replay finished ones —", 10, False, None),
    ("  this is especially useful for demo day when the advisor wants to observe a session", 10, False, None),
    "",
    ("Developer documentation created:", 13, True, None),
    ("• DESIGN.md — UI design system (colors, spacing, z-index, animations)", 11, False, None),
    ("• BOOT.md — Setup, run, test, debug guide", 11, False, None),
    ("• CONTRIBUTING.md — Branch naming, PR process, code conventions", 11, False, None),
    ("• ARCHITECTURE.md — Data flow, AI models, session store, known gaps", 11, False, None),
])


# ── References — original slide 5 ──
slide5 = prs.slides[5]
for shape in slide5.shapes:
    if shape.name == "Content Placeholder 2":
        clear_and_set(shape, [
            ("[1] N. Montfort, Twisty Little Passages: An Approach to Interactive Fiction. MIT Press, 2003.",
             12, False, None),
            ("[2] A. Plotkin, The Inform Designer's Handbook, 2004.",
             12, False, None),
            ("[3] P. Ammanabrolu et al., How to Avoid Being Eaten by a Grue, arXiv:2006.07409, 2020.",
             12, False, None),
            ("[4] N. Walton, AI Dungeon, Latitude, 2019. https://aidungeon.io",
             12, False, None),
            ("[5] OpenAI, GPT-5 API Documentation, 2025. https://platform.openai.com/docs",
             12, False, None),
            ("[6] N. Shaker et al., Procedural Content Generation in Games. Springer, 2016.",
             12, False, None),
            ("[7] Socket.IO Documentation, 2024. https://socket.io/docs/v4/",
             12, False, None),
            ("[8] WebRTC API, MDN Web Docs, 2024. https://developer.mozilla.org/en-US/docs/Web/API/WebRTC_API",
             12, False, None),
        ])


# ── Slide 6: Attachment ────────────────────────────────────────
slide6 = prs.slides[6]
for shape in slide6.shapes:
    if shape.name == "Content Placeholder 2":
        clear_and_set(shape, [
            ("Live Application:", 16, True, None),
            ("• https://velvet-shadow.web.app", 14, False, None),
            ("• Backend: Cloud Run (europe-west1)", 14, False, None),
            ("• Database: Firestore (comp491-akboys-2026)", 14, False, None),
            "",
            ("Submitted Documents:", 16, True, None),
            ("• Registration Form, Proposal (original + updated)", 14, False, None),
            ("• Progress Meeting 1, 2, 3 presentations", 14, False, None),
            "",
            ("GitHub Repository:", 16, True, None),
            ("• https://github.com/bakaraman/comp491-akboys", 14, False, None),
            ("• 8 issues closed, 6 PRs merged, 27 open issues for WP5", 14, False, None),
        ], font_size=14)


# ── Reorder slides: move added slides (7,8,9) after slide 4 ──
# Current order: 0-Title, 1-Info, 2-Prev, 3-Current, 4-Kadir, 5-Refs, 6-Attach, 7-Serdar, 8-Ata, 9-Batuhan
# Desired order: 0-Title, 1-Info, 2-Prev, 3-Current, 4-Kadir, 5-Serdar, 6-Ata, 7-Batuhan, 8-Refs, 9-Attach
slide_list = prs.slides._sldIdLst
slides = list(slide_list)
# Desired order by current index: [0, 1, 2, 3, 4, 7, 8, 9, 5, 6]
desired = [0, 1, 2, 3, 4, 7, 8, 9, 5, 6]
reordered = [slides[i] for i in desired]
for s in slides:
    slide_list.remove(s)
for s in reordered:
    slide_list.append(s)

# ── Save ───────────────────────────────────────────────────────
out = "/Users/batuhankaraman/comp491-akboys/docs/filled/COMP491_Progress3.pptx"
prs.save(out)
print(f"Saved: {out}")
