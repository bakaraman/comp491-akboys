from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation("../lecture-slides/Comp491_ProgressPresentationTemplate_fixed.pptx")


def clear_and_set(shape, lines, font_size=18, bold_first=False):
    """Replace all paragraphs in a text frame with new lines."""
    tf = shape.text_frame
    for i in range(len(tf.paragraphs) - 1, 0, -1):
        p = tf.paragraphs[i]._p
        p.getparent().remove(p)
    for li, line in enumerate(lines):
        if li == 0:
            para = tf.paragraphs[0]
            para.clear()
        else:
            para = tf.add_paragraph()

        if isinstance(line, tuple):
            text, sz, bld, color = line
        else:
            text = line
            sz = font_size
            bld = (li == 0 and bold_first)
            color = None

        run = para.add_run()
        run.text = text
        run.font.size = Pt(sz)
        run.font.bold = bld
        if color:
            run.font.color.rgb = RGBColor(*color)
        para.space_after = Pt(3)
        para.space_before = Pt(0)


# ── Slide 0: Title Slide ───────────────────────────────────────
slide0 = prs.slides[0]
txBox = slide0.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(8), Inches(2.5))
tf = txBox.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.space_after = Pt(12)
r = p.add_run()
r.text = "Text-Based Adventure with LLMs"
r.font.size = Pt(32)
r.font.bold = True
r.font.color.rgb = RGBColor(255, 255, 255)

p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
p2.space_after = Pt(8)
r2 = p2.add_run()
r2.text = "COMP 491 — Progress Meeting 2"
r2.font.size = Pt(20)
r2.font.color.rgb = RGBColor(220, 220, 220)

p3 = tf.add_paragraph()
p3.alignment = PP_ALIGN.CENTER
r3 = p3.add_run()
r3.text = "March 12, 2026"
r3.font.size = Pt(16)
r3.font.color.rgb = RGBColor(200, 200, 200)


# ── Slide 1: Project Info ──────────────────────────────────────
slide1 = prs.slides[1]
for shape in slide1.shapes:
    if shape.name == "Content Placeholder 2":
        clear_and_set(shape, [
            ("Project Title: Text-Based Adventure with LLMs", 18, True, None),
            "",
            ("Team Members:", 16, True, None),
            ("  Kadir Yigit Ozcelik (79975)", 15, False, None),
            ("  Serdar Yengil (80232)", 15, False, None),
            ("  Batuhan Karaman (79791)", 15, False, None),
            ("  Ata Berke Goktekin (80277)", 15, False, None),
            "",
            ("Advisor: Baris Akgun", 16, True, None),
            ("Progress Meeting 2 — March 12, 2026", 15, False, None),
        ], font_size=15)


# ── Slide 2: Previous Progress Meeting ─────────────────────────
slide2 = prs.slides[2]
for shape in slide2.shapes:
    if shape.name == "Content Placeholder 2":
        clear_and_set(shape, [
            ("Progress Meeting 1 Summary (March 5, 2026):", 16, True, None),
            "",
            ("What was accomplished:", 15, True, None),
            ("• GitHub repo created with branch protection and issue tracking", 13, False, None),
            ("• Single-player demo built (Tkinter + Gemini API, 5 rooms, 3 NPCs)", 13, False, None),
            ("• Architecture documented (Structured Chaos: referee + LLM narrator)", 13, False, None),
            ("• Project proposal and registration form submitted", 13, False, None),
            "",
            ("What was planned next:", 15, True, None),
            ("• Design scenario JSON schema and validation rules", 13, False, None),
            ("• Implement LLM-based scenario generation", 13, False, None),
            ("• Start universe templates (noir, cyberpunk, etc.)", 13, False, None),
            ("• Begin backend migration from Python demo", 13, False, None),
        ], font_size=13)


# ── Slide 3: Current Progress Meeting ──────────────────────────
slide3 = prs.slides[3]
for shape in slide3.shapes:
    if shape.name == "Content Placeholder 2":
        clear_and_set(shape, [
            ("WP2 — Core Engine, Scenarios & Web Platform (W3-W6)", 15, True, None),
            "",
            ("Major Accomplishments:", 14, True, None),
            ("• Full migration from Python/Tkinter to TypeScript monorepo", 12, False, None),
            ("  (Node.js + Express backend, Next.js 15 + React 19 frontend)", 12, False, None),
            ("• 6 complete game scenarios designed and implemented:", 12, False, None),
            ("  Noir, Haunted Manor, Space Station, Pirate, Western, Cyberpunk", 12, False, None),
            ("  Each with 5 rooms, 3 NPCs (with dialogue), 5 items (evidence system)", 12, False, None),
            ("• OpenAI GPT-5.4 streaming narration via Server-Sent Events (SSE)", 12, False, None),
            ("• GPT-5-nano integration for real-time follow-up action suggestions", 12, False, None),
            ("• Session management with UUID-based routing (/session/[id])", 12, False, None),
            ("• Rich markdown rendering (bold, italic, blockquotes, headings)", 12, False, None),
            ("• Dark noir-themed responsive UI with scenario selection grid", 12, False, None),
            ("• REST API with 6 endpoints (scenarios, sessions, chat, suggestions)", 12, False, None),
            "",
            ("Work Distribution:", 14, True, None),
            ("• Batuhan: Monorepo setup, GCP infra, server architecture, API design", 12, False, None),
            ("• Kadir: Next.js frontend, scenario picker UI, chat components", 12, False, None),
            ("• Serdar: Express backend, session store, SSE streaming", 12, False, None),
            ("• Ata Berke: 6 scenario scripts (rooms, NPCs, items), prompt engineering", 12, False, None),
        ], font_size=12)


# ── Slide 4: Next Progress Meeting ─────────────────────────────
slide4 = prs.slides[4]
for shape in slide4.shapes:
    if shape.name == "Content Placeholder 2":
        clear_and_set(shape, [
            ("Next: WP3 — Multiplayer, Persistence & Advanced Features (W7-W10)", 15, True, None),
            "",
            ("Planned Tasks:", 14, True, None),
            ("• Implement explicit GameState tracking on server (room, inventory, visited)", 12, False, None),
            ("• Add Socket.IO for real-time multiplayer co-op sessions", 12, False, None),
            ("• Migrate from in-memory SessionStore to Firestore for persistence", 12, False, None),
            ("• Deploy backend to Google Cloud Run, frontend to Firebase Hosting", 12, False, None),
            ("• Add image generation for scene visualization (gpt-image-1.5)", 12, False, None),
            ("• Implement game-over conditions and win/lose scenarios", 12, False, None),
            ("• Add user authentication (Firebase Auth)", 12, False, None),
            "",
            ("Work Distribution:", 14, True, None),
            ("• Batuhan: GCP deployment (Cloud Run + Firebase), Firestore migration", 12, False, None),
            ("• Kadir: Multiplayer UI, player indicators, shared game view", 12, False, None),
            ("• Serdar: Socket.IO backend, multiplayer session management", 12, False, None),
            ("• Ata Berke: Image generation prompts, game-over logic, balancing", 12, False, None),
        ], font_size=12)


# ── Slide 5: References ────────────────────────────────────────
slide5 = prs.slides[5]
for shape in slide5.shapes:
    if shape.name == "Content Placeholder 2":
        clear_and_set(shape, [
            ("[1] N. Montfort, Twisty Little Passages: An Approach to Interactive Fiction. MIT Press, 2003.",
             12, False, None),
            ("[2] A. Plotkin, The Inform Designer's Handbook, 2004.",
             12, False, None),
            ("[3] P. Ammanabrolu et al., How to Avoid Being Eaten by a Grue, arXiv:2006.07409, 2020.",
             12, False, None),
            ("[4] N. Walton, AI Dungeon, Latitude, 2019. https://aidungeon.io",
             12, False, None),
            ("[5] OpenAI, GPT-5 API Documentation, 2025. https://platform.openai.com/docs",
             12, False, None),
            ("[6] N. Shaker et al., Procedural Content Generation in Games. Springer, 2016.",
             12, False, None),
            ("[7] Socket.IO Documentation, 2024. https://socket.io/docs/v4/",
             12, False, None),
        ])


# ── Slide 6: Attachment ────────────────────────────────────────
slide6 = prs.slides[6]
for shape in slide6.shapes:
    if shape.name == "Content Placeholder 2":
        clear_and_set(shape, [
            ("Submitted Documents:", 16, True, None),
            ("• COMP491_Registration_Form_FILLED.docx", 14, False, None),
            ("• COMP491_Proposal_FILLED.docx (with Gantt chart)", 14, False, None),
            ("• COMP491_Progress1.pptx (Progress Meeting 1)", 14, False, None),
            "",
            ("Working Application:", 16, True, None),
            ("• Full-stack web app: Next.js + Express + OpenAI GPT-5.4", 14, False, None),
            ("• 6 playable scenarios with streaming AI narration", 14, False, None),
            ("• Live demo available at http://localhost:3000", 14, False, None),
            "",
            ("GitHub Repository:", 16, True, None),
            ("• https://github.com/bakaraman/comp491-akboys", 14, False, None),
        ], font_size=14)


# ── Save ───────────────────────────────────────────────────────
out = "/Users/batuhankaraman/comp491-akboys/docs/filled/COMP491_Progress2.pptx"
prs.save(out)
print(f"Saved: {out}")
