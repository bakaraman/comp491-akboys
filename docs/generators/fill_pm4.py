"""
fill_pm4.py — Generate COMP491_Progress4.pptx (and .pdf) plus PM4_Presentation_Script.pdf

- Uses Comp491_ProgressPresentationTemplate_fixed.pptx for the slide deck.
- Uses textboxes (NOT the bullet-styled content placeholder) so we control
  the bullet character ourselves; this removes the double-bullet visual bug
  ('▪ • ...') that appeared when content already started with '•' AND the
  layout master added its own bullet.
- Reportlab generates the matching speaker-script PDF.

Run:
  cd docs/generators && python3 fill_pm4.py
Outputs:
  ../../COMP491_Progress4.pptx
  ../../COMP491_Progress4.pdf            (via soffice)
  ../../PM4_Presentation_Script.pdf
"""

from pathlib import Path
import subprocess

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib.colors import HexColor
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.platypus.flowables import KeepTogether
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont


def _register_unicode_font() -> str:
    """Register a TTF font that supports Turkish characters (ş, ğ, ı, İ, ç, ö, ü).

    reportlab's built-in Helvetica is PDF Type 1 and lacks Turkish glyphs, so
    Turkish chars render as ■ boxes. Falls back across a few system locations.
    """
    candidates = [
        ("/System/Library/Fonts/Supplemental/Arial.ttf",
         "/System/Library/Fonts/Supplemental/Arial Bold.ttf"),
        ("/Library/Fonts/Arial.ttf",
         "/Library/Fonts/Arial Bold.ttf"),
        ("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
         "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"),
    ]
    for regular, bold in candidates:
        if Path(regular).exists() and Path(bold).exists():
            pdfmetrics.registerFont(TTFont("UIFont", regular))
            pdfmetrics.registerFont(TTFont("UIFont-Bold", bold))
            pdfmetrics.registerFontFamily("UIFont", normal="UIFont", bold="UIFont-Bold")
            return "UIFont"
    return "Helvetica"


REPO_ROOT = Path(__file__).resolve().parents[2]
TEMPLATE = REPO_ROOT / "docs" / "lecture-slides" / "Comp491_ProgressPresentationTemplate_fixed.pptx"

OUT_PPTX = REPO_ROOT / "COMP491_Progress4.pptx"
OUT_PDF = REPO_ROOT / "COMP491_Progress4.pdf"
OUT_SCRIPT_PDF = REPO_ROOT / "PM4_Presentation_Script.pdf"


# ──────────────────────────────────────────────────────────────────────
#  CONTENT
# ──────────────────────────────────────────────────────────────────────

TITLE = "Text-Based Adventure with LLMs"
COURSE_LINE = "COMP 491 — Progress Meeting 4"
DATE_LINE = "May 8, 2026"
ADVISOR = "Baris Akgun"
TEAM = [
    ("Kadir Yigit Ozcelik", "79975"),
    ("Serdar Yengil", "80232"),
    ("Batuhan Karaman", "79791"),
    ("Ata Berke Goktekin", "80277"),
]

PM3_DONE = [
    "Socket.IO multiplayer with room codes, lobby, action batching",
    "Firebase Auth (Google sign-in) and server-side GameState tracking",
    "Structured directive validation (MOVE, PICKUP, UNLOCK)",
    "gpt-image-1.5 scene generation with per-scenario style prompts",
    "Accusation mechanic with evidence chain validation",
    "Backend on Google Cloud Run, frontend on Firebase Hosting",
    "Firestore session persistence",
]
PM3_PLAN = [
    "Game architecture pivot to host-prompt procedural worlds",
    "Multiplayer game-over flow, turn limits, finale cinematic",
    "Voice chat, evidence/notes panel, in-game tutorial",
    "Spectator mode and replay system",
    "Procedural narrative depth (red herrings, alibis, NPC backstories)",
    "Operational hardening (timeouts, usage logging)",
]

CURRENT_OVERVIEW = [
    ("WP5 — Procedural Worlds, Voice Chat, Spectator Mode, Narrative Depth", True),
    ("The game pivoted from six hardcoded scenarios to a fully procedural pipeline driven by a host prompt; the codebase ships demo-ready on production.", False),
    ("", False),
    ("Velvet Shadow v2 — the new game model:", True),
    ("The host types a theme; the LLM generates the world (rooms, NPCs, items, the murder, the solution) in roughly thirty seconds.", False),
    ("Single-player removed from the UI; multiplayer-only, two to ten detectives per session.", False),
    ("Opening and finale are cinematic — AI image, Turkish TTS narration, ambient score, synced typewriter reveal.", False),
    ("", False),
    ("Kadir — Frontend, audio, voice and onboarding:", True),
    ("Three-layer ambient soundtrack and global UI click sound, with a settings popover for separate music and SFX volume.", False),
    ("Post-game crime-scene reconstruction modal — six to eight chronological beats walking the team through what really happened.", False),
    ("Push-to-talk voice chat over WebRTC mesh, plus a shared Evidence Board notepad that replaces the old chat tabs.", False),
    ("Five-step in-game tutorial overlay shown automatically on first play, replayable from the settings menu.", False),
    ("", False),
    ("Serdar — Spectators, replay and backend hardening:", True),
    ("New spectator role at the socket layer — read-only access to a live game, with all write handlers rejecting their input.", False),
    ("Replay page for finished sessions: scrubbable timeline, play / pause and 1x / 2x / 5x speed, ending with the reconstruction reveal.", False),
    ("Zod payload validation on every client-emitted event, plus an integration test that verifies Firestore persistence survives a restart.", False),
    ("", False),
    ("Ata — Procedural narrative depth and genre awareness:", True),
    ("Every NPC now ships with an alibi object and a short backstory; the culprit's alibi must contain a discoverable inconsistency.", False),
    ("Each world contains at least one red herring item that the narrator presents with the weight of real evidence.", False),
    ("Genre detection on the host prompt injects a tonal style guide — noir, sci-fi, gothic-horror, cyberpunk, period drama.", False),
    ("", False),
    ("Batuhan — The pivot itself, AI integration and production deploy (see slide 5)", True),
]

# Per-person slide content (heading bold flag, body bullets)
BATUHAN_BULLETS_DONE = [
    "Led the Velvet Shadow v2 pivot — replaced six hardcoded scenarios with an LLM-generated world pipeline (strict Zod schema, structured outputs, an automatic repair pass and a hardcoded fallback world) so the host types a theme and the game produces its own rooms, NPCs, evidence and culprit.",
    "Removed roughly two hundred lines of regex movement and pickup heuristics and collapsed the entire directive surface to a single MOVE — inventory, evidence, sanity and NPC mood now live in narrator prose, giving the model full authorial control.",
    "Cut the pre-game wait from sixty seconds to thirty by running opening-image generation in parallel with world-text generation, and resolved the resulting race so an early image is baked straight into the ready event instead of being dropped by an old client filter.",
    "Rewrote the TTS pipeline — moved off broken WAV streams that some browsers refused to decode, switched to MP3, swapped to the warmer 'ash' voice and pre-fetched the audio at story-ready, so the opening cinematic plays at zero latency.",
    "Shipped real-time multiplayer broadcasts: players in the same room now see each other's actions and the narrator stream live, a session-wide queue banner shows who is currently being written to, and a fifteen-second disconnect grace period keeps reloads and tab throttling from flickering players off the roster.",
    "Ran today's full production deploy — backend on Cloud Run in europe-west1, web on Cloud Run in us-central1, environment baked at build time, Firebase Hosting CDN auto-invalidated, and the live site verified at velvet-shadow.web.app.",
]
BATUHAN_BULLETS_NEXT = [
    "Add OpenAI-client timeouts, exponential-backoff retry and a fallback model cascade so a single 5xx no longer terminates a live session.",
    "Add per-IP and per-session rate limiting plus usage and cost telemetry to Cloud Logging so a runaway client cannot blow up the bill.",
    "Integrate Sentry and a small admin observability dashboard to surface which prompts hit 429s and which sessions end unexpectedly.",
]

KADIR_BULLETS_DONE = [
    "Built three independent ambient audio layers — lobby, in-game and cinematic — with cross-fade ducking, plus a settings popover that gives separate volume sliders and mute toggles for music and UI sound effects, with preferences persisted to localStorage.",
    "Added a global UI click sound that hooks every button through document-level event delegation, so the audio touch is consistent across the whole app.",
    "Built the post-game crime-scene reconstruction modal that opens after the finale and walks the player through six to eight chronological beats; strict ID enums keep the AI from inventing rooms or NPCs, and a single-flight cache prevents duplicate generations across browsers.",
    "Split reconstruction generation into two model calls — events on a reasoning model, conclusion on a fast non-reasoning one — after the original single call kept truncating the closing paragraph against the same token budget.",
    "Shipped V-key push-to-talk voice chat over WebRTC mesh with Socket.IO acting as a signaling-only relay; holding V transmits to all teammates, typing 'v' inside any input field passes through cleanly, and the speaker's avatar lights up live.",
    "Replaced the old room and direct chat tabs with a single shared Evidence Board notepad — every note broadcasts to all players, treated as out-of-game thinking material that the narrator never reads.",
    "Built a five-step in-game tutorial overlay that triggers on first play and highlights the turn counter, chat area, action input, accuse button and settings; keyboard navigable and replayable from the settings menu.",
]
KADIR_BULLETS_NEXT = [
    "Polish the mobile and small-screen layout — header buttons, the accusation banner and the action input still need responsive media queries.",
    "Run end-to-end voice-chat verification against the production HTTPS deployment, where getUserMedia and STUN traversal behave differently than on localhost.",
    "Add an in-app changelog panel for the demo audience so newcomers can see what shipped in the latest sprint.",
]

SERDAR_BULLETS_DONE = [
    "Added a spectator role to the Socket.IO handshake so observers receive the full session state and the live narrator stream while every write handler rejects them — invited guests can watch a real game without being able to interfere.",
    "Built the spectator UI on the home page and adapted the in-game view for it: chat input, accuse and leave are hidden and a Spectator badge appears in the header.",
    "Built the replay system as a dedicated page with a scrubbable timeline, play / pause and 1x / 2x / 5x speed controls, ending with the same reconstruction modal as the live game.",
    "Added Zod schema validation to every client-emitted Socket.IO event so malformed or hostile payloads never reach a handler — the server returns an INVALID_PAYLOAD error event instead.",
    "Wrote a Firestore integration test that creates a session, populates it, syncs, simulates a server restart, rehydrates the store and asserts byte-for-byte equality of the recovered state.",
]
SERDAR_BULLETS_NEXT = [
    "Add session-resume on the spectator side so a dropped observer reconnects and catches up automatically without a manual refresh.",
    "Profile and tune the action-batching window under realistic four-player load, since the current values were chosen against solo tests.",
    "Add a server-side per-session rate limit on action submission to complement the per-IP limit at the OpenAI call layer.",
]

ATA_BULLETS_DONE = [
    "Extended the procedural world schema with three new fields — every NPC now ships with an alibi object (claimed location, claimed activity, optional corroboration and inconsistency) and a short backstory the narrator can pull from.",
    "Made every world contain at least one red-herring item, and constrained the generator so the culprit's alibi must hide a discoverable inconsistency while innocent NPCs stay consistent but information-limited.",
    "Updated the narrator prompt to use the new fields so when a player questions a suspect the narrator delivers their alibi naturally, draws on the backstory for voice and gives red herrings the same narrative weight as real evidence.",
    "Updated the post-game reconstruction so the conclusion explicitly acknowledges any red herring (\"the broken letter was a misdirection\"), letting players see which leads were intentional dead ends.",
    "Built a genre-aware style amplifier — a keyword heuristic with a small LLM fallback detects the genre of the host prompt and injects a tonal style guide into both the world generator and the narrator's system prompt.",
    "Authored concrete style guides for noir, sci-fi, gothic-horror, cyberpunk and period drama, plus a generic fallback — the same procedural pipeline now produces tonally distinct worlds across themes.",
]
ATA_BULLETS_NEXT = [
    "Broaden the genre library with fantasy, modern thriller and cosmic horror, and tighten the noir guide based on playtest feedback.",
    "Add a pre-game world-quality check where the LLM critiques its own world and re-rolls if it spots unreachable rooms, missing inconsistencies or no red herring.",
    "Run a five-prompt diversity playtest to measure whether the new procedural depth genuinely shifts gameplay between themes.",
]

REFERENCES = [
    "[1] N. Montfort, Twisty Little Passages: An Approach to Interactive Fiction. MIT Press, 2003.",
    "[2] A. Plotkin, The Inform Designer's Handbook, 2004.",
    "[3] P. Ammanabrolu et al., How to Avoid Being Eaten by a Grue, arXiv:2006.07409, 2020.",
    "[4] N. Walton, AI Dungeon, Latitude, 2019. https://aidungeon.io",
    "[5] OpenAI, GPT-5 API Documentation, 2025. https://platform.openai.com/docs",
    "[6] N. Shaker et al., Procedural Content Generation in Games. Springer, 2016.",
    "[7] Socket.IO Documentation, 2024. https://socket.io/docs/v4/",
    "[8] WebRTC API, MDN Web Docs, 2024. https://developer.mozilla.org/en-US/docs/Web/API/WebRTC_API",
]

ATTACHMENT_LINES = [
    ("Live Application:", True),
    ("https://velvet-shadow.web.app", False),
    ("Backend: Cloud Run (europe-west1)", False),
    ("Database: Firestore (comp491-akboys-2026)", False),
    ("", False),
    ("GitHub Repository:", True),
    ("https://github.com/bakaraman/comp491-akboys", False),
    ("8 issues closed this sprint, 2 PRs merged + 1 branch direct-merged", False),
    ("2 open issues for WP6 (Batuhan's operational hardening — in progress)", False),
]


# ──────────────────────────────────────────────────────────────────────
#  PPTX generation
# ──────────────────────────────────────────────────────────────────────

RED = RGBColor(0xC8, 0x10, 0x2E)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
LIGHT = RGBColor(0xFF, 0xFF, 0xFF)


def _wipe(shape):
    tf = shape.text_frame
    # Drop every paragraph except the first; clear that one
    for i in range(len(tf.paragraphs) - 1, 0, -1):
        p = tf.paragraphs[i]._p
        p.getparent().remove(p)
    tf.paragraphs[0].clear()
    return tf


def _set_para(para, text, *, size=14, bold=False, color=DARK, align=None, bullet=False, space_after=4):
    if align is not None:
        para.alignment = align
    para.space_after = Pt(space_after)
    para.space_before = Pt(0)
    if bullet:
        text = f"•  {text}"
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_textbox(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tb.text_frame.word_wrap = True
    return tb


def _add_section_title(slide, text):
    """Add a black page title at the top with a thin red rule beneath it."""
    tb = _add_textbox(slide, 0.5, 0.35, 9.0, 0.8)
    tf = _wipe(tb)
    p = tf.paragraphs[0]
    _set_para(p, text, size=28, bold=True, color=DARK)
    # red rule
    rule = slide.shapes.add_shape(1, Inches(0.5), Inches(1.15), Inches(9.0), Inches(0.04))
    rule.fill.solid()
    rule.fill.fore_color.rgb = RED
    rule.line.fill.background()


def _content_box(slide, *, top=1.35, height=5.9):
    # Slide is 10" wide × 7.5" tall. Leave 0.5" margins on left/right.
    return _add_textbox(slide, 0.5, top, 9.0, height)


def _bullet(tf, text, *, size=14, bold=False, color=DARK, space_after=4):
    if not tf.paragraphs[0].runs and not tf.paragraphs[0].text:
        para = tf.paragraphs[0]
    else:
        para = tf.add_paragraph()
    _set_para(para, text, size=size, bold=bold, color=color, bullet=True, space_after=space_after)
    return para


def _line(tf, text, *, size=14, bold=False, color=DARK, blank_above=False, space_after=4):
    if blank_above:
        sp = tf.add_paragraph()
        _set_para(sp, "", size=size)
    if not tf.paragraphs[0].runs and not tf.paragraphs[0].text:
        para = tf.paragraphs[0]
    else:
        para = tf.add_paragraph()
    _set_para(para, text, size=size, bold=bold, color=color, space_after=space_after)
    return para


def build_pptx() -> Path:
    prs = Presentation(str(TEMPLATE))

    # Reset to a clean deck — keep only the title slide layout (template's slide 0)
    title_layout = prs.slides[0].slide_layout
    blank_layout = next(l for l in prs.slide_layouts if l.name == "Title Only")

    # Remove every default slide and rebuild
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')  # noqa
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    # ── Slide 0: title ──
    # The title-slide layout is a blank red background with the KOÇ UNIVERSITY
    # mark baked in; it has no placeholders. We overlay a centered textbox
    # towards the bottom half so it doesn't collide with the logo (top half).
    s = prs.slides.add_slide(title_layout)
    tb = _add_textbox(s, 0.5, 5.7, 9.0, 1.7)
    tf = _wipe(tb)
    _set_para(tf.paragraphs[0], TITLE,
              size=28, bold=True, color=LIGHT, align=PP_ALIGN.CENTER)
    _set_para(tf.add_paragraph(), COURSE_LINE,
              size=16, color=LIGHT, align=PP_ALIGN.CENTER)
    _set_para(tf.add_paragraph(), DATE_LINE,
              size=14, color=LIGHT, align=PP_ALIGN.CENTER)

    # ── Slide 1: project info ──
    s = prs.slides.add_slide(blank_layout)
    _add_section_title(s, "Comp491 Project Progress")
    tb = _content_box(s)
    tf = tb.text_frame
    tf.paragraphs[0].clear()
    _set_para(tf.paragraphs[0], f"Project Title: {TITLE}", size=18, bold=True, color=RED)
    _line(tf, "Team Members:", size=18, bold=True, color=RED, blank_above=True)
    for name, sid in TEAM:
        _bullet(tf, f"{name} ({sid})", size=15)
    _line(tf, f"Advisor: {ADVISOR}", size=18, bold=True, color=RED, blank_above=True)
    _line(tf, f"{COURSE_LINE} — {DATE_LINE}", size=14)

    # ── Slide 2: Previous Progress Meeting (PM3 summary) ──
    s = prs.slides.add_slide(blank_layout)
    _add_section_title(s, "Previous Progress Meeting")
    tb = _content_box(s, top=1.30, height=6.05)
    tf = tb.text_frame
    tf.paragraphs[0].clear()
    _set_para(tf.paragraphs[0], "Progress Meeting 3 Summary (April 3, 2026):",
              size=14, bold=True, color=RED, space_after=4)
    _line(tf, "What was accomplished in WP3+WP4:", size=12, bold=True, color=RED, blank_above=True, space_after=2)
    for b in PM3_DONE:
        _bullet(tf, b, size=11, space_after=3)
    _line(tf, "What was planned for WP5:", size=12, bold=True, color=RED, blank_above=True, space_after=2)
    for b in PM3_PLAN:
        _bullet(tf, b, size=11, space_after=3)

    # ── Slide 3: Current Progress Meeting (overview) ──
    s = prs.slides.add_slide(blank_layout)
    _add_section_title(s, "Current Progress Meeting")
    tb = _content_box(s, top=1.30, height=6.05)
    tf = tb.text_frame
    tf.paragraphs[0].clear()
    first = True
    for text, is_heading in CURRENT_OVERVIEW:
        if text == "":
            sp = tf.add_paragraph()
            _set_para(sp, "", size=4, space_after=0)
            continue
        if is_heading:
            if first:
                _set_para(tf.paragraphs[0], text, size=11, bold=True, color=RED, space_after=2)
                first = False
            else:
                _line(tf, text, size=11, bold=True, color=RED, space_after=2)
        else:
            _bullet(tf, text, size=9, space_after=2)

    # ── Slide 4: Batuhan ──
    _person_slide(
        prs, blank_layout,
        title="Batuhan Karaman — Infrastructure & GCP",
        completed_label="Completed (Architecture pivot, AI optimization, production deploy):",
        completed=BATUHAN_BULLETS_DONE,
        next_up=BATUHAN_BULLETS_NEXT,
    )

    # ── Slide 5: Kadir ──
    _person_slide(
        prs, blank_layout,
        title="Kadir Yigit Ozcelik — Frontend & UI",
        completed_label="Completed (Issues #24, #36, #48, #49):",
        completed=KADIR_BULLETS_DONE,
        next_up=KADIR_BULLETS_NEXT,
    )

    # ── Slide 6: Serdar ──
    _person_slide(
        prs, blank_layout,
        title="Serdar Yengil — Backend & API",
        completed_label="Completed (Issues #52, #53):",
        completed=SERDAR_BULLETS_DONE,
        next_up=SERDAR_BULLETS_NEXT,
    )

    # ── Slide 7: Ata ──
    _person_slide(
        prs, blank_layout,
        title="Ata Berke Goktekin — Game Design & Prompts",
        completed_label="Completed (Issues #54, #55):",
        completed=ATA_BULLETS_DONE,
        next_up=ATA_BULLETS_NEXT,
    )

    # ── Slide 8: References ──
    s = prs.slides.add_slide(blank_layout)
    _add_section_title(s, "References")
    tb = _content_box(s)
    tf = tb.text_frame
    tf.paragraphs[0].clear()
    first = True
    for ref in REFERENCES:
        if first:
            _set_para(tf.paragraphs[0], ref, size=12, color=RED)
            first = False
        else:
            _line(tf, ref, size=12, color=RED)

    # ── Slide 9: Attachment ──
    s = prs.slides.add_slide(blank_layout)
    _add_section_title(s, "Attachment")
    tb = _content_box(s)
    tf = tb.text_frame
    tf.paragraphs[0].clear()
    first = True
    for text, is_heading in ATTACHMENT_LINES:
        if text == "":
            sp = tf.add_paragraph()
            _set_para(sp, "", size=10)
            continue
        if is_heading:
            if first:
                _set_para(tf.paragraphs[0], text, size=16, bold=True, color=RED)
                first = False
            else:
                _line(tf, text, size=16, bold=True, color=RED)
        else:
            _bullet(tf, text, size=13, color=RED)

    OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PPTX))
    return OUT_PPTX


def _person_slide(prs, layout, *, title, completed_label, completed, next_up):
    s = prs.slides.add_slide(layout)
    _add_section_title(s, title)
    tb = _content_box(s, top=1.30, height=6.05)
    tf = tb.text_frame
    tf.paragraphs[0].clear()

    # Heuristic font sizing: more bullets → smaller font so a single slide holds them all.
    n = len(completed) + len(next_up)
    if n >= 10:
        bullet_size, bullet_gap, header_size = 10, 3, 13
    elif n >= 9:
        bullet_size, bullet_gap, header_size = 11, 3, 14
    else:
        bullet_size, bullet_gap, header_size = 12, 4, 14

    _set_para(tf.paragraphs[0], completed_label, size=header_size, bold=True, color=RED, space_after=3)
    for b in completed:
        _bullet(tf, b, size=bullet_size, space_after=bullet_gap)
    _line(tf, "Next Up:", size=header_size, bold=True, color=RED, blank_above=True, space_after=3)
    for b in next_up:
        _bullet(tf, b, size=bullet_size, space_after=bullet_gap)
    return s


def pptx_to_pdf(pptx_path: Path, pdf_path: Path) -> None:
    """Convert pptx → pdf via LibreOffice (soffice). Outputs into pdf_path."""
    out_dir = pdf_path.parent
    cmd = [
        "soffice",
        "--headless",
        "--convert-to",
        "pdf",
        "--outdir",
        str(out_dir),
        str(pptx_path),
    ]
    subprocess.run(cmd, check=True)
    # soffice writes <stem>.pdf in out_dir; rename if needed
    produced = out_dir / (pptx_path.stem + ".pdf")
    if produced != pdf_path:
        produced.replace(pdf_path)


# ──────────────────────────────────────────────────────────────────────
#  Speaker-script PDF (Turkish)
# ──────────────────────────────────────────────────────────────────────

# Turkish speaker script — written to be read aloud, in full sentences, no jargon dump.
SCRIPT_TR = {
    "Batuhan Karaman": {
        "subtitle": "Mimari Pivot, Yapay Zeka Entegrasyonu ve Yayın",
        "done": [
            "Bu sprint'te oyunun temelini değiştirdim. PM3'te elimizde altı tane elle yazılmış senaryo vardı; host bunlardan birini seçiyordu. Velvet Shadow v2 ile bunu tamamen kaldırdım: artık host bir kutuya istediği temayı yazıyor — mesela 1927 Chicago'da kayıp bir caz şarkıcısı — ve yapay zeka otuz saniye içinde kendi başına komple bir noir gizem dünyası kuruyor; odalar, NPC'ler, kanıtlar ve cinayetin çözümü dahil her şey procedural üretiliyor.",
            "Bu pivotun mümkün olabilmesi için oyun motorunda yaklaşık iki yüz satırlık eski regex tabanlı hareket ve eşya alma kodunu sildim ve bütün directive yüzeyini tek bir MOVE komutuna indirdim. Eşya almak, kanıt bulmak, NPC ruh hali değişimleri — bunların hepsini artık kod değil, anlatıcı doğal proza olarak ifade ediyor. Hikâye üzerindeki tam yetkiyi modele verdim, oyun çok daha akıcı bir hâle geldi.",
            "Oyun başlamadan önce gösterilen sinematik açılışı ciddi şekilde hızlandırdım. Önceden açılış görseli ile dünya yazısı sırayla üretiliyordu, toplam bekleme yaklaşık altmış saniyeydi. İkisini paralel hale getirdim ve aralarındaki yarış durumunu — görsel önce biterse hangi event'le payload'a girdiğini — server tarafında çözdüm. Şu an bekleme süresi otuz saniye.",
            "Sesli anlatımı baştan yazdım. Eski WAV formatı bazı tarayıcılarda bozuk çalıyordu, MP3'e geçtim, modelin sesini daha sıcak ve audiobook tonundaki 'ash' sesine çevirdim, ve dünya hazır olur olmaz sesi önceden indirmeye başlattım — böylece kullanıcı 'Perdeyi Aç'a bastığında ses anında oynamaya başlıyor, beklemeden.",
            "Çok-oyunculu tarafta gerçek-zaman katmanını oturttum. Aynı odadaki oyuncular artık birbirinin yazdığını ve anlatıcının kendilerine verdiği cevabı canlı görebiliyor; ekranın üstünde tüm oyunculara 'anlatıcı şu an kime yazıyor, sırada kim var' diye gösteren bir banner var. Bağlantı koparsa on beş saniyelik bir bekleme süresi koydum, böylece sayfa yenileme veya kısa bir kesinti oyuncuyu listeden düşürmüyor.",
            "Son olarak bugün, sunum sabahı, tüm sistemi production'a deploy ettim. Backend europe-west1'de Cloud Run'da çalışıyor, frontend us-central1'de yine Cloud Run üstünde, ortam değişkenleri Docker build sırasında bundle'a gömülüyor ve Firebase Hosting CDN her sürümde otomatik invalidate ediliyor. Canlı sürüm velvet-shadow.web.app adresinde, az önce uçtan uca test ettim, çalışıyor.",
        ],
        "next": [
            "OpenAI istemcisine zaman aşımı, exponential backoff ile yumuşak yeniden deneme ve yedek model kademesi ekleyeceğim — şu an tek bir 5xx canlı oyunu düşürebiliyor.",
            "IP başına ve oturum başına rate limit ile birlikte Cloud Logging'e kullanım ve maliyet telemetrisi yazacağım — çünkü tek bir kullanıcının aksiyon spam'i şu an faturayı patlatabilir.",
            "Sentry'yi entegre edip küçük bir admin gözlemleme paneli yapacağım — Cloud Logging tek başına hangi prompt'un 429 yediğini veya hangi oturumun beklenmedik şekilde kapandığını anlamak için yeterince ince taneli değil.",
        ],
    },
    "Kadir Yigit Ozcelik": {
        "subtitle": "Frontend, Ses, Voice Chat ve Onboarding",
        "done": [
            "Oyuna sıfırdan ses kattım. Üç bağımsız atmosfer katmanı yaptım — ana menü ve lobi için bir döngü, oyun içi için ayrı bir döngü, sinematikler için zaten var olan üçüncü katman — aralarında yumuşak geçiş var. Ayarlar paneline müzik ve UI ses efektleri için ayrı ayrı ses sürgüleri ve sustur düğmeleri ekledim, tercihler tarayıcıda saklanıyor.",
            "Tüm butonlara global bir tıklama sesi bağladım. Bunu doküman seviyesinde yapan bir delegasyon mantığıyla yazdım, böylece her yeni butonu manuel sarmaya gerek yok; arayüzün tutarlı bir dokunuşu oldu.",
            "Oyun finali bittikten sonra açılan 'Olay Yerini Yeniden Canlandır' modalını yaptım. Modal altı ila sekiz kronolojik sahnede gerçekte ne olduğunu sırayla anlatıyor; her sahne gerçek bir oda ve gerçek bir NPC'yi gösteriyor. Modeli sıkı ID listeleriyle bağladım, sahte oda veya karakter uyduramıyor; aynı oturuma ikinci bir tarayıcıdan giren olursa tek bir generation'ı paylaşıyorlar.",
            "Bu reconstruction'ı yazarken üretimi iki ayrı modele böldüm — olay zinciri için reasoning yapan bir model, sonuç paragrafı için hızlı bir non-reasoning model — çünkü tek modelde sonuç token bütçesinin son kelimesinde kesiliyordu, ayırınca temizlendi.",
            "V tuşuna basılı tutarak konuşulan sesli sohbet ekledim. WebRTC mesh ile peer-to-peer ses akıyor, Socket.IO sadece signaling yapıyor. Bir input alanına yazarken V tuşu ses göndermiyor — yani 'evet' yazarken yanlışlıkla yayın yapılmıyor — ve konuşan oyuncunun avatarında yeşil halka beliriyor.",
            "Eski 'Bu Oda / Direkt' iletişim panelini tek paylaşımlı bir 'Kanıt Tahtası'na dönüştürdüm. Her oyuncu kanıt teorilerini buraya yazıyor, herkes anında görüyor; anlatıcı buradakileri okumuyor, yani oyuncuların oyun-dışı düşünce alanı haline geldi.",
            "Oyuna ilk giren kullanıcı için beş adımlık bir tutorial overlay'i ekledim. Tur sayacını, sohbet alanını, aksiyon kutusunu, suçlama düğmesini ve ayarları sırayla vurguluyor; klavyeyle gezinilebilir ve ayarlar menüsünden tekrar açılabiliyor.",
        ],
        "next": [
            "Mobil ve küçük ekran düzenlemesini bitireceğim — başlık butonları, suçlama bandı ve aksiyon kutusu hâlâ uygun media query'lere ihtiyaç duyuyor.",
            "Sesli sohbeti production HTTPS sürümünde uçtan uca test edeceğim — getUserMedia ve STUN traversal localhost dışında farklı davranıyor, gerçek deploy'da doğrulamadan rahat olamam.",
            "Demo izleyicisi için bir 'Bu Sürümde Yenilikler' paneli ekleyeceğim — son sprint'te o kadar çok şey eklendi ki, ilk açan kullanıcının bunları görmesi gerekiyor.",
        ],
    },
    "Serdar Yengil": {
        "subtitle": "İzleyici Modu, Replay ve Backend Sağlamlaştırma",
        "done": [
            "Socket.IO el sıkışmasına bir 'izleyici' rolü ekledim. Demo gününde danışmanımız ve diğer izleyiciler devam eden bir oyunu canlı izleyebilecek ama oyuna müdahale edemeyecek; izleyici soketi tüm oturum durumunu ve canlı anlatıcı akışını alıyor, ancak her oyuncu-aksiyon handler'ı izleyiciyi reddediyor.",
            "İzleyici arayüzünü ana sayfaya yerleştirdim — oda kodu kutusunun altında 'İzleyici olarak katıl' düğmesi var. Oyun ekranında izleyici için sohbet kutusu, suçlama ve ayrılma butonlarını gizledim, başlığa 'İzleyici' rozeti ekledim, böylece sahnelendiği rol baştan belli oluyor.",
            "Bitmiş oyunlar için ayrı bir replay sayfası yaptım. Sayfanın üstünde sürüklenebilir bir zaman çizelgesi var, oynat/duraklat ve 1x / 2x / 5x hız kontrolleri çalışıyor; sürgü ilerledikçe mesajlar sırayla açığa çıkıyor ve sayfa Kadir'in yaptığı reconstruction modalıyla bitiyor.",
            "İstemcinin gönderdiği her Socket.IO event'ine Zod şema doğrulaması ekledim. Bozuk veya kötü niyetli payload artık handler'a ulaşmıyor; sunucu doğrudan INVALID_PAYLOAD hata event'iyle yanıt veriyor, böylece API yüzeyi sertleşti.",
            "Firestore oturum deposu için bir entegrasyon testi yazdım. Test bir oturum yaratıyor, içine mesajlar ve dünya verisi koyuyor, sync ediyor, sunucu yeniden başlatma simülasyonu yapıyor, rehydrate ediyor ve geri gelen state'in byte-for-byte eşit olduğunu doğruluyor — yani sunucu çökse de tek bir oyun verisi kaybedilmiyor.",
        ],
        "next": [
            "İzleyici tarafına oturum-devam etme özelliği ekleyeceğim — bağlantısı kopan bir izleyici manuel yenileme yapmadan kaldığı yerden devam edebilsin.",
            "Aksiyon batching penceresini gerçekçi dört oyunculu yük altında profilleyip ayarlayacağım — şu anki değerler tek oyunculu testte seçilmişti, kalabalık oyunda davranışı doğrulamadık.",
            "Oyuncu aksiyon gönderimi için sunucu tarafında oturum başına bir rate limit ekleyeceğim — Batuhan'ın OpenAI çağrı seviyesindeki IP limiti ile beraber iki katmanlı koruma olacak.",
        ],
    },
    "Ata Berke Goktekin": {
        "subtitle": "Procedural Anlatı Derinliği ve Tür Bilinci",
        "done": [
            "Procedural dünya şemasına üç yeni alan ekledim. Artık her NPC bir alibi nesnesi taşıyor — nerede olduğunu iddia ettiği yer, ne yaptığı, gerekirse onu doğrulayan başka bir karakter ve sadece katiller için: alibinin gizli tutarsızlığı. Her NPC'nin ayrıca anlatıcının ses tonu için kullanabileceği iki üç cümlelik bir geçmiş hikâyesi var.",
            "Üretici prompt'una iki kural koydum: her dünya mutlaka en az bir kırmızı şaşırtmaca içerecek, ve katilin alibisi mutlaka oyuncunun bulabileceği somut bir tutarsızlık taşıyacak. Masum NPC'lerin alibileri tutarlı ama bilgi açısından sınırlı, böylece sorgulama gerçek bir dedektif iş gibi hissettiriyor.",
            "Anlatıcı prompt'unu yeni alanları kullanacak şekilde güncelledim. Oyuncu bir NPC'yi sıkıştırdığında anlatıcı o karakterin alibisini doğal bir dille anlatıyor, geçmişinden gelen sesi kullanıyor; kırmızı şaşırtmacalar gerçek kanıtlar kadar ağırlıkla sahnelenip oyuncunun kafasını karıştırıyor.",
            "Reconstruction prompt'unu da güncelledim — sonuç paragrafı artık varsa kırmızı şaşırtmacalara değiniyor, mesela 'kırık mektup bir yanıltmaydı' gibi. Böylece oyuncu oyun bittiğinde hangi izlerin kasıtlı çıkmaz olduğunu da öğreniyor.",
            "Tür bilinçli bir stil yükseltici modül yaptım. Host prompt'tan türü tespit ediyorum — önce anahtar kelime sezgisi, gerekirse küçük bir LLM çağrısı — sonra türe özgü stil rehberini hem dünya üreticisinin hem de anlatıcının sistem prompt'una enjekte ediyorum.",
            "Beş tür için somut stil rehberi yazdım: noir, bilim kurgu, gotik korku, siberpunk ve dönem draması. Aynı procedural boru hattı artık aynı temayla farklı dünyalarda tonal olarak ciddi biçimde ayrışıyor — sci-fi prompt'u yazdığında atmosfer gerçekten metalik ve klinik geliyor.",
        ],
        "next": [
            "Stil kütüphanesini fantastik, modern gerilim ve kozmik korku ekleyerek genişleteceğim, ve playtest geri bildirimine göre noir rehberini sıkılaştıracağım.",
            "Üretim öncesi bir kalite doğrulama adımı ekleyeceğim — LLM kendi ürettiği dünyayı eleştiriyor ve ulaşılamayan oda, çelişkisi olmayan alibi veya eksik kırmızı şaşırtmaca tespit ederse yeniden üretiyor.",
            "Yeni procedural derinliğin temalar arasında gerçekten farklı oyun deneyimi yarattığını ölçmek için beş prompt'luk bir çeşitlilik playtest'i koşacağım.",
        ],
    },
}


def build_script_pdf() -> Path:
    doc = SimpleDocTemplate(
        str(OUT_SCRIPT_PDF), pagesize=A4,
        leftMargin=0.9 * inch, rightMargin=0.9 * inch,
        topMargin=0.85 * inch, bottomMargin=0.85 * inch,
        title="PM4 Sunum Metni",
    )

    styles = getSampleStyleSheet()
    base_font = _register_unicode_font()
    title_style = ParagraphStyle(
        "title", parent=styles["Title"], fontName=base_font,
        fontSize=24, leading=30, spaceAfter=12, textColor=HexColor("#111111"),
    )
    sub_style = ParagraphStyle(
        "sub", parent=styles["BodyText"], fontName=base_font,
        fontSize=12, leading=16, textColor=HexColor("#666666"),
        alignment=1,
    )
    h1 = ParagraphStyle(
        "h1", parent=styles["Heading1"], fontName=base_font,
        fontSize=24, leading=28, spaceAfter=4, textColor=HexColor("#111111"),
    )
    h_sub = ParagraphStyle(
        "h_sub", parent=styles["BodyText"], fontName=base_font,
        fontSize=11, leading=14, textColor=HexColor("#888888"), spaceAfter=18,
    )
    h2 = ParagraphStyle(
        "h2", parent=styles["Heading2"], fontName=base_font,
        fontSize=14, leading=18, spaceBefore=12, spaceAfter=6, textColor=HexColor("#111111"),
    )
    bullet_style = ParagraphStyle(
        "bullet", parent=styles["BodyText"], fontName=base_font,
        fontSize=10.2, leading=14, leftIndent=14, bulletIndent=2,
        spaceAfter=4, textColor=HexColor("#222222"),
    )
    intro_style = ParagraphStyle(
        "intro", parent=styles["BodyText"], fontName=base_font,
        fontSize=11.5, leading=16, alignment=1, textColor=HexColor("#333333"),
    )

    story = []

    # Cover page
    story.append(Spacer(1, 1.4 * inch))
    story.append(Paragraph("Progress Meeting 4 — Sunum Metni", title_style))
    story.append(Paragraph("Text-Based Adventure with LLMs", sub_style))
    story.append(Paragraph("COMP 491 — 8 Mayıs 2026", sub_style))
    story.append(Spacer(1, 0.6 * inch))
    story.append(Paragraph(
        "Bu doküman her takım üyesinin PM4 sunumunda söyleyeceği metni içerir.<br/>"
        "Her kişi için ayrı sayfa: yaptıkları ve yapacakları, Türkçe, bullet point.",
        intro_style,
    ))

    # Order: Batuhan first (per his request), then alphabetical-ish
    order = ["Batuhan Karaman", "Kadir Yigit Ozcelik", "Serdar Yengil", "Ata Berke Goktekin"]

    for name in order:
        story.append(PageBreak())
        block = SCRIPT_TR[name]
        story.append(Paragraph(name, h1))
        story.append(Paragraph(block["subtitle"], h_sub))

        story.append(Paragraph("Yaptıklarım", h2))
        for b in block["done"]:
            story.append(Paragraph(b, bullet_style, bulletText="•"))

        story.append(Paragraph("Yapacaklarım", h2))
        for b in block["next"]:
            story.append(Paragraph(b, bullet_style, bulletText="•"))

    doc.build(story)
    return OUT_SCRIPT_PDF


# ──────────────────────────────────────────────────────────────────────
#  main
# ──────────────────────────────────────────────────────────────────────

def main() -> None:
    print("[fill_pm4] generating pptx …")
    pptx = build_pptx()
    print(f"  → {pptx.relative_to(REPO_ROOT)}")

    print("[fill_pm4] converting pptx → pdf via soffice …")
    pptx_to_pdf(pptx, OUT_PDF)
    print(f"  → {OUT_PDF.relative_to(REPO_ROOT)}")

    print("[fill_pm4] generating speaker-script pdf …")
    script_pdf = build_script_pdf()
    print(f"  → {script_pdf.relative_to(REPO_ROOT)}")

    print("[fill_pm4] done.")


if __name__ == "__main__":
    main()
